"""填充 Smart-Cat 演示 PPT 的脚本。
基于 docs/模板ppt.pptx 的四部分结构，将演示PPT大纲.md 的内容填入对应文本框。
策略：递归遍历所有形状（含组合），保留模板的字体/颜色/位置，只替换文字。
"""
from pptx import Presentation
from pptx.util import Pt
import copy

SRC = "d:/vscode/Smart-Cat/docs/Smart-Cat演示.pptx"
DST = "d:/vscode/Smart-Cat/docs/Smart-Cat演示.pptx"


def iter_shapes(shapes):
    """递归遍历所有形状，包括组合内的子形状。"""
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:  # GROUP
            yield from iter_shapes(shape.shapes)


def set_text_preserve_format(shape, new_text):
    """替换形状文字，保留首个 run 的格式。"""
    if not shape.has_text_frame:
        return False
    tf = shape.text_frame
    # 保留第一个 run 的格式
    first_para = tf.paragraphs[0] if tf.paragraphs else None
    first_run = first_para.runs[0] if first_para and first_para.runs else None

    # 清空所有段落
    for para in tf.paragraphs:
        for run in list(para.runs):
            run.text = ""

    if first_run:
        first_run.text = new_text
    else:
        # 没有 run，直接设置 text
        tf.text = new_text
    return True


def find_and_replace(slide, target_name, new_text):
    """按 name 找到形状并替换文字。"""
    for shape in iter_shapes(slide.shapes):
        if shape.name == target_name:
            return set_text_preserve_format(shape, new_text)
    return False


def replace_by_current_text(slide, old_text, new_text):
    """按当前文字内容匹配并替换（用于无 name 或 name 重复的情况）。"""
    for shape in iter_shapes(slide.shapes):
        if shape.has_text_frame and shape.text_frame.text.strip() == old_text.strip():
            return set_text_preserve_format(shape, new_text)
    return False


def main():
    prs = Presentation(SRC)

    # ===== 封面 (slide 0) =====
    s = prs.slides[0]
    # 主标题文本框：TextBox 113 含英文副标题 -> 改为项目定位
    find_and_replace(s, 'TextBox 113',
        'Smart-Cat 边缘智能称重识别一体机')
    # 答辩人信息
    find_and_replace(s, '矩形 11', '答辩人：Smart-Cat 项目组      时间：2026年7月')
    # 副标题装饰
    find_and_replace(s, 'TextBox 113',
        'Smart-Cat 边缘智能称重识别一体机')
    # 创意小标签改为主标语
    replace_by_current_text(s, '创意', '边缘 AI')

    # ===== 目录 (slide 1) =====
    s = prs.slides[1]
    replace_by_current_text(s, '选题意义', '项目背景与目标')
    replace_by_current_text(s, '本文论点', '系统架构与流程')
    replace_by_current_text(s, '本文结构', '识别与云边协同')
    replace_by_current_text(s, '本文结语', '演示与总结')

    # ===== Part 01 章节封面 (slide 2) =====
    s = prs.slides[2]
    replace_by_current_text(s, '选题意义', '项目背景与目标')
    replace_by_current_text(s, 'THE TITLE CONTENT', 'BACKGROUND & GOALS')

    # ===== slide 3: 选题背景概述 -> 项目背景与问题 =====
    s = prs.slides[3]
    replace_by_current_text(s, '选题背景概述', '为什么要做智能称重识别')
    # 主体说明文字
    replace_by_current_text(s,
        '您的内容打在这里，或者通过复制您的文本后，在此框中选择粘贴，并选择只保留文字您的内容打在这里，或者通过复制您的文本后，在此框中选择粘贴，并选择只保留文字',
        '传统称重收银依赖人工识别商品，单价、重量、总价容易录入错误；交易数据和设备状态缺少统一管理；高峰场景下操作效率低。本系统希望把识别、称重、计价、播报自动化，并由 Web 后台统一管理。')

    # ===== slide 4: 选题意义概述 -> 项目整体目标 =====
    s = prs.slides[4]
    replace_by_current_text(s, '选题意义概述', '从单点识别到完整业务闭环')

    # ===== slide 5: 数据展示页 -> 硬件构成 =====
    s = prs.slides[5]
    replace_by_current_text(s, '选题意义概述', '硬件构成与系统定位')
    # 4 个百分比文本框 -> 硬件关键词
    replace_by_current_text(s, '94%', 'RK3576')
    replace_by_current_text(s, '42%', 'HX711')
    replace_by_current_text(s, '23%', 'SYN6288')
    replace_by_current_text(s, '63%', 'YOLO11')
    # 底部说明
    replace_by_current_text(s,
        '您的内容打在这里，或者通过复制您的文本后，在此框中选择粘贴，并选择只保留文字您的内容打在这里，或者通过复制您的文本后，在此框中选择粘贴，并选择只保留文字',
        '鲁班猫 3 / RK3576 内置 NPU，外接 USB 摄像头、HX711 称重模块、SYN6288 语音合成模块与 USB 无线网卡，构成边缘智能称重识别一体机。')

    # ===== slide 6: 选题发展方向 -> 系统主流程 =====
    s = prs.slides[6]
    replace_by_current_text(s, '选题发展方向', '系统主流程：从放商品到后台同步')
    replace_by_current_text(s, '个人素质', '完整闭环')
    replace_by_current_text(s,
        '在持续长期发展过程中，继承优良传统，适应时代要求，由企业家积极倡导，全体员工自觉实践，从而形成的代表企业信念、激发企业活力、推动企业生产经营的团体精神和行为规范。按照不同层次划分为精神文化型和组织制度型两类。划分为精神文化型和组织制度型两类',
        '放上商品 → 称重触发 → 摄像头拍照 → RKNN/NPU 商品检测 → HX711 读取重量 → 查询单价并计算总价 → SYN6288 语音播报 → JSONL 本地记录 → MQTT 上报 → Web 后台同步显示。强调：不是单纯跑模型，而是完整业务闭环。')

    # ===== Part 02 章节封面 (slide 7) =====
    s = prs.slides[7]
    replace_by_current_text(s, '本文论点', '系统架构与流程')
    replace_by_current_text(s, 'THE TITLE CONTENT', 'ARCHITECTURE & FLOW')

    # ===== slide 8: 选题整体分析 -> 系统架构四层 =====
    s = prs.slides[8]
    replace_by_current_text(s, '选题整体分析', '系统架构：边缘端 + MQTT + Web 后台')
    titles = ['添加您的标题', '您的内容打在这里，或者通过复制您的文本后，在此框中选择粘贴，并选择只保留文字。']
    arch_content = [
        ('感知层', 'USB 摄像头、HX711 称重模块、称重传感器、SYN6288 语音合成模块'),
        ('边缘计算层', '鲁班猫 3 / RK3576，负责拍照、称重、NPU 推理、计价、播报'),
        ('通信协同层', 'MQTT 负责交易上报、设备事件、运行参数、策略下发和语音命令'),
        ('平台管理层', '电脑 Web 后台负责交易流水、商品维护、统计分析、设备状态和策略管理'),
    ]
    # 4 个 Text Box 10，按顺序替换标题与内容
    box_idx = 0
    for shape in iter_shapes(s.shapes):
        if shape.name == 'Text Box 10' and shape.has_text_frame:
            if box_idx < 4:
                title, content = arch_content[box_idx]
                set_text_preserve_format(shape, f'{title}：{content}')
                box_idx += 1

    # ===== slide 9: 选题发展过程 -> 板端常驻交易流程 =====
    s = prs.slides[9]
    replace_by_current_text(s, '选题发展过程', '板端常驻称重触发流程')

    # ===== slide 10: 选题存在问题 -> AI 模型与 RKNN/NPU 部署 =====
    s = prs.slides[10]
    replace_by_current_text(s, '选题存在问题', '为什么选择 RKNN / NPU 边缘推理')
    replace_by_current_text(s,
        '单击此处添加文字，添加具体文字内容，您的说明文字在此处添加，此处输入详细的说明文字。',
        '训练阶段 YOLO11 检测模型 → 导出 ONNX → 转换 RKNN → RK3576 NPU 推理。RK3576 NPU 单张约 30-80ms，远优于树莓派 4B 的约 3000ms，让边缘设备本地实时识别，不依赖云端。')

    # ===== Part 03 章节封面 (slide 11) =====
    s = prs.slides[11]
    replace_by_current_text(s, '本文结构', '识别与云边协同')
    replace_by_current_text(s, 'THE TITLE CONTENT', 'RECOGNITION & CLOUD-EDGE')

    # ===== slide 12: 选题研究思路 -> 识别状态与人工确认机制 =====
    s = prs.slides[12]
    replace_by_current_text(s, '选题研究思路', '不是所有识别结果都直接结算')
    replace_by_current_text(s, 'SUCCESS', '待确认')

    # ===== slide 13: 选题数据分析 -> 语音补盲与无障碍查询 =====
    s = prs.slides[13]
    replace_by_current_text(s, '选题数据分析', '语音补盲：不只修正商品，也能查询状态')
    replace_by_current_text(s,
        '您的内容打在这里，或者通过复制您的文本后，在此框中选择粘贴，并选择只保留字您的您的内容打在这里，或者通过复制您的文本后，在此框中选择粘贴，并选择只保留字您的您的内容打在这里，或者通过复制您的文本后，在此框中选择粘贴，并选择只保留字您的',
        '视障顾客说“重量”→ TTS 念当前克数；说“价格”→ 念单价总价；说“最近交易”→ 念上一笔结果；店员说“状态”→ 念设备运行情况。SYN6288 不只是报价输出，而是整个人机交互入口。')

    # ===== slide 14: 选题数据分析(2) -> 云边协同与策略下发 =====
    s = prs.slides[14]
    replace_by_current_text(s, '选题数据分析', '相同模型 + 不同策略 = 不同业务效果')
    # 7 个"添加标题"，分别替换为策略相关要点
    policy_points = [
        '模型负责识别：画面里像什么商品',
        '策略负责决策：业务上如何处理',
        '宽松策略：可直接结算',
        '保守策略：进入待确认',
        '促销策略：按促销规则计价',
        '后台 MQTT retained 下发新策略',
        '板端下一笔交易即按新规则执行',
    ]
    idx = 0
    for shape in iter_shapes(s.shapes):
        if shape.name == 'TextBox 692' and shape.has_text_frame and idx < len(policy_points):
            set_text_preserve_format(shape, policy_points[idx])
            idx += 1

    # ===== slide 15: 选题运用研究 -> MQTT 消息设计 =====
    s = prs.slides[15]
    replace_by_current_text(s, '选题运用研究', 'MQTT：普通消息与 Retained 消息的分工')
    replace_by_current_text(s,
        '您的内容打在这里，或者通过复制您的文本后，在此框中选择粘贴，并选择只保留文字您的内容打在这里，您的内容打在这里或者通过',
        '交易记录、设备事件、语音命令用普通消息（一次性事件）；运行参数、设备策略用 Retained 消息（新设备上线立即拿到最新规则）。普通消息适合事件，retained 适合状态。')

    # ===== slide 16: 选题运用研究(2) -> Web 后台功能 =====
    s = prs.slides[16]
    replace_by_current_text(s, '选题运用研究', 'Web 后台：管理、统计、配置与补盲')
    replace_by_current_text(s, '添加您的标题', '后台可远程改变边缘设备行为')
    replace_by_current_text(s,
        '单击此处添加文字，添加具体文字内容，您的说明文字在此处添加，此处输入详细的说明文字。单击此处添加文字',
        '交易流水：筛选、详情、CSV 导出、清空流水。商品管理：名称、单价、单位、启用状态、价格历史。')
    replace_by_current_text(s, '添加具体文字内容，您的说明文字在此处添加，此处输入详细的说明文字。',
        '运行参数与设备策略支持 MQTT retained 下发，低置信度交易可由手机语音补盲修正。')

    # ===== slide 17: 选题案例分析 -> Web 后台页面清单 =====
    s = prs.slides[17]
    replace_by_current_text(s, '选题案例分析', 'Web 后台五大页面')
    replace_by_current_text(s, '添加您的标题', '监控概览：交易数、销售额、总重量、异常/待确认')
    replace_by_current_text(s,
        '您的内容打在这里，或者通过复制您的文本后，在此框中选择粘贴，并选择只保留文字。您的内容打在这里，您的内容打在这里，或者通过复制您的文本后，在此框中选择粘贴您的内容打在这里，或者通过复制您的文本后，在此框中选择粘贴，并选择只保留文字。您的内容',
        '运行参数：置信度阈值、语音音量、MQTT 开关。设备策略：低置信度动作、计价模式、语音模板、回滚。设备事件：心跳、服务状态、策略应用、语音命令记录。语音补盲：手机输入“这是苹果”等修正交易。')

    # ===== Part 04 章节封面 (slide 18) =====
    s = prs.slides[18]
    replace_by_current_text(s, '本文结语', '演示与总结')
    replace_by_current_text(s, 'THE TITLE CONTENT', 'DEMO & SUMMARY')

    # ===== slide 19: 选题研究总结 -> 现场演示流程 =====
    s = prs.slides[19]
    replace_by_current_text(s, '选题研究总结', '现场演示流程')

    # ===== slide 20: 选题未来应用分析 -> 总结与改进方向 =====
    s = prs.slides[20]
    replace_by_current_text(s, '选题未来应用分析', '项目完成度与改进方向')
    summary_points = [
        ('硬件闭环已完成', '摄像头 + HX711 + SYN6288 三件套联动'),
        ('边缘 AI 推理已完成', 'YOLO11 → ONNX → RKNN → NPU 本地实时识别'),
        ('云边协同已完成', 'MQTT 参数/策略下发，事件回传'),
        ('后续改进方向', 'JSONL 升级 SQLite；多设备管理；权限完善；外壳优化'),
    ]
    idx = 0
    for shape in iter_shapes(s.shapes):
        if shape.name == 'Text Box 10' and shape.has_text_frame and idx < len(summary_points):
            title, content = summary_points[idx]
            set_text_preserve_format(shape, f'{title}：{content}')
            idx += 1

    # ===== slide 21: 选题发展规划 -> 项目价值总结 =====
    s = prs.slides[21]
    replace_by_current_text(s, '选题发展规划', '项目价值总结')
    replace_by_current_text(s, '添加您的内容.', '边缘端独立完成交易，不依赖云端识别')
    # 其余三个 "添加您的内容" 用 set 顺序替换
    value_points = [
        'MQTT 让后台与板端协同，可远程配置',
        '低置信度有人工兜底，不是假自动化',
        'AI 模型真正嵌入可运行的业务闭环',
    ]
    idx = 0
    for shape in iter_shapes(s.shapes):
        if shape.name.startswith('矩形') and shape.has_text_frame and shape.text_frame.text.strip() == '添加您的内容' and idx < len(value_points):
            set_text_preserve_format(shape, value_points[idx])
            idx += 1

    # ===== slide 22: 文献综述 -> 删除（保留空白或改致谢延伸） =====
    s = prs.slides[22]
    replace_by_current_text(s, '选题文献综述', '技术资料与项目文档')
    replace_by_current_text(s, '行业PPT模板http://www.1ppt.com/hangye/',
        '详细文档见 docs/ 目录：项目说明、API接口、商业计划书、开发日志')

    # ===== slide 23: 谢谢聆听 =====
    s = prs.slides[23]
    replace_by_current_text(s, '答辩人：第一PPT      时间：2028年12月31日',
        'Smart-Cat 边缘智能称重识别一体机      2026年7月')
    replace_by_current_text(s, 'BY:第一PPT', 'BY: Smart-Cat')

    prs.save(DST)
    print(f"已保存：{DST}")


if __name__ == "__main__":
    main()
