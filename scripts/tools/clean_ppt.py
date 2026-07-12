"""补充清理 Smart-Cat 演示 PPT 中的残留占位文字。
针对 build_ppt.py 运行后仍残留的模板装饰文字做统一替换或清空。
"""
from pptx import Presentation

PATH = "d:/vscode/Smart-Cat/docs/Smart-Cat演示.pptx"


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:  # GROUP
            yield from iter_shapes(shape.shapes)


def set_text_keep_fmt(shape, new_text):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    first_para = tf.paragraphs[0] if tf.paragraphs else None
    first_run = first_para.runs[0] if first_para and first_para.runs else None
    for para in tf.paragraphs:
        for run in list(para.runs):
            run.text = ""
    if first_run:
        first_run.text = new_text
    else:
        tf.text = new_text


def clean_slide(slide_idx, replacements):
    """对指定 slide 执行多组 (旧文字片段, 新文字) 替换。
    replacements 中旧文字用 in 判断包含关系。"""
    prs = Presentation(PATH)
    s = prs.slides[slide_idx]
    for shape in iter_shapes(s.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        for old_frag, new_text in replacements:
            if old_frag in text:
                set_text_keep_fmt(shape, new_text)
                break
    prs.save(PATH)


def main():
    # slide 3: 清理残留装饰文字
    clean_slide(3, [
        ("选题背景概述", "为什么要做智能称重识别"),
        ("此部分内容作为文字排版占位", "人工识别慢、易出错、数据割裂，是称重收银三大痛点"),
        ("25%", "效率低"),
        ("51%", "易出错"),
    ])

    # slide 4: 清理重复的"选题意义概述"
    clean_slide(4, [
        ("此部分内容作为文字排版占位", "设备自动触发交易，本地识别称重计价，MQTT 上报，后台统一管理"),
    ])

    # slide 5: 清理
    clean_slide(5, [
        ("此部分内容作为文字排版占位", "边缘 NPU 本地推理，无需云端，低延迟实时识别"),
    ])

    # slide 6: 清理残留的"执行力/创新力/感召力"和长段落
    clean_slide(6, [
        ("执行力", "触发"),
        ("创新力", "识别"),
        ("感召力", "播报"),
        ("在持续长期发展过程中", "放上商品后自动触发称重，NPU 识别商品，HX711 读重量，查价计价，SYN6288 播报，JSONL 记录，MQTT 上报后台，形成完整闭环"),
    ])

    # slide 8: SWOT 图 - 替换为架构四层关键词
    clean_slide(8, [
        ("优势", "感知层"),
        ("劣势", "计算层"),
        ("机会", "协同层"),
        ("威胁", "管理层"),
    ])

    # slide 9: "添加您的标题" -> 板端流程步骤
    prs = Presentation(PATH)
    s = prs.slides[9]
    flow_steps = [
        "空秤待机，预加载 RKNN，保持摄像头打开",
        "重量超过阈值，多次采样确认稳定后拍照",
        "NPU 推理识别商品，HX711 稳定重量参与计价",
        "SYN6288 播报报价，写入记录并 MQTT 上报",
    ]
    idx = 0
    for shape in iter_shapes(s.shapes):
        if shape.has_text_frame and "添加您的标题" in shape.text_frame.text and idx < len(flow_steps):
            set_text_keep_fmt(shape, flow_steps[idx])
            idx += 1
    prs.save(PATH)

    # slide 10: 清理残留的"添加您的标题"和占位文字
    clean_slide(10, [
        ("单击此处添加文字，添加具体文字内容", "RK3576 NPU 单张约 30-80ms，树莓派 4B 约 3000ms，PC CPU 约 500-1000ms。NPU 让边缘设备本地实时识别"),
    ])
    # 清理剩余"添加您的标题"
    prs = Presentation(PATH)
    s = prs.slides[10]
    rknn_points = ["训练 YOLO11", "导出 ONNX", "转换 RKNN", "NPU 推理"]
    idx = 0
    for shape in iter_shapes(s.shapes):
        if shape.has_text_frame and "添加您的标题" in shape.text_frame.text and idx < len(rknn_points):
            set_text_keep_fmt(shape, rknn_points[idx])
            idx += 1
    prs.save(PATH)

    # slide 12: 清理"文字内容/请替换文字内容"
    clean_slide(12, [
        ("请替换文字内容，点击添加相关标题文字", "Top1 置信度 >= 0.75 且 Top1-Top2 差值 >= 0.15 时自动结算，否则进入待确认，不计算总价，等待人工修正"),
        ("文字内容", "自动结算条件"),
    ])

    # slide 13: 清理百分比和占位
    clean_slide(13, [
        ("80%", "重量"),
        ("95%", "价格"),
        ("75%", "最近交易"),
        ("您的内容打在这里，或者通过复制您的文本后，在此框中选择粘贴，并选择只保留字您的", "视障顾客说“重量”→TTS 念当前克数；说“价格”→念单价总价；说“最近交易”→念上一笔结果"),
    ])

    # slide 15: 清理百分比和占位
    clean_slide(15, [
        ("80%", "交易记录"),
        ("您的内容打在这里，或者通过复制您的文本后，在此框中选择粘贴，并选择只保留字您的", "交易、事件、语音命令用普通消息；参数、策略用 Retained 消息，新设备上线即取最新规则"),
    ])

    # slide 16: 清理残留占位
    clean_slide(16, [
        ("单击此处添加文字，添加具体文字内容", "交易流水筛选/详情/CSV 导出；商品管理含价格历史；统计分析含销量排行与识别质量"),
    ])

    # slide 17: 清理"案例一/案例二"和占位
    clean_slide(17, [
        ("案例一", "监控概览"),
        ("案例二", "交易流水"),
        ("案例三", "商品管理"),
        ("您的内容打在这里，或者通过复制您的文本后", "后台可远程下发参数与策略，低置信度交易可由手机语音补盲修正，形成可管理的数据闭环"),
    ])

    # slide 19: 演示流程 - 替换"添加您的标题"
    prs = Presentation(PATH)
    s = prs.slides[19]
    demo_steps = [
        "启动 Web 后台与板端常驻称重服务",
        "放上商品，自动识别称重计价，SYN6288 播报",
        "Web 后台同步交易流水，展示设备事件与统计",
        "低置信度交易用手机语音补盲修正，下发新策略",
    ]
    idx = 0
    for shape in iter_shapes(s.shapes):
        if shape.has_text_frame:
            t = shape.text_frame.text
            if ("添加您的标题" in t or "您的内容打在这里" in t) and idx < len(demo_steps):
                set_text_keep_fmt(shape, demo_steps[idx])
                idx += 1
    prs.save(PATH)

    # slide 20: SWOT -> 完成度
    clean_slide(20, [
        ("优势", "已完成"),
        ("劣势", "待改进"),
        ("机会", "扩展方向"),
        ("威胁", "技术风险"),
    ])
    # 替换 SWOT 细节
    prs = Presentation(PATH)
    s = prs.slides[20]
    done_points = [
        "硬件闭环：摄像头+HX711+SYN6288",
        "边缘推理：YOLO11→ONNX→RKNN→NPU",
        "云边协同：MQTT 参数/策略下发",
        "后续：SQLite、多设备、权限完善",
    ]
    idx = 0
    for shape in iter_shapes(s.shapes):
        if shape.has_text_frame and ("添加您的标题" in shape.text_frame.text or "您的内容打在这里" in shape.text_frame.text) and idx < len(done_points):
            set_text_keep_fmt(shape, done_points[idx])
            idx += 1
    prs.save(PATH)

    # slide 21: 清理剩余"添加您的内容"
    prs = Presentation(PATH)
    s = prs.slides[21]
    for shape in iter_shapes(s.shapes):
        if shape.has_text_frame and shape.text_frame.text.strip() == "添加您的内容":
            set_text_keep_fmt(shape, "可运行、可配置、可管理的业务闭环")
    prs.save(PATH)

    # slide 22: 清理"文献一/二/三"和占位
    clean_slide(22, [
        ("文献一", "项目说明"),
        ("文献二", "API 接口文档"),
        ("文献三", "商业计划书"),
        ("您的内容打在这里，或者通过复制您的文本后", "详细文档见 docs/ 目录，包含项目说明、API接口、商业计划书、开发日志、脚本参数说明"),
    ])

    print("清理完成")


if __name__ == "__main__":
    main()
