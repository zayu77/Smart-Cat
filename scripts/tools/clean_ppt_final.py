"""最终清理 Smart-Cat 演示 PPT 中所有残留的模板占位文字。
基于 verify_ppt.py 的检查结果，精确处理每一处残留。
"""
from pptx import Presentation

PATH = "d:/vscode/Smart-Cat/docs/Smart-Cat演示.pptx"


def iter_shapes(shapes):
    """递归遍历所有形状，包括组合内的子形状。"""
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:  # GROUP
            yield from iter_shapes(shape.shapes)


def set_text_keep_fmt(shape, new_text):
    """替换形状文字，保留首个 run 的格式。"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    first_para = tf.paragraphs[0] if tf.paragraphs else None
    first_run = first_para.runs[0] if first_para and first_para.runs else None
    for para in tf.paragraphs:
        for run in list(para.runs):
            run.text = ""
    if first_run:
        first_run.text = new_text
    else:
        tf.text = new_text


def contains_any(text, keywords):
    """检查文本是否包含任一关键词。"""
    return any(kw in text for kw in keywords)


def clean_slide(slide_idx, replacements):
    """对指定 slide 执行多组 (关键词列表, 新文字) 替换。"""
    prs = Presentation(PATH)
    s = prs.slides[slide_idx]
    for shape in iter_shapes(s.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        for keywords, new_text in replacements:
            if contains_any(text, keywords):
                set_text_keep_fmt(shape, new_text)
                break
    prs.save(PATH)


def replace_sequential(slide_idx, match_keywords, new_texts):
    """按顺序替换匹配的形状文字（处理重复占位符）。"""
    prs = Presentation(PATH)
    s = prs.slides[slide_idx]
    idx = 0
    for shape in iter_shapes(s.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if idx < len(new_texts) and contains_any(text, match_keywords):
            set_text_keep_fmt(shape, new_texts[idx])
            idx += 1
    prs.save(PATH)
    return idx


def main():
    # slide 1: 目录页残留的 "D", "D", "esign" 装饰文字 -> 清空无用碎片
    # 这些是 "DESIGN" 被拆分后的装饰文字，保留主目录文本即可
    prs = Presentation(PATH)
    s = prs.slides[1]
    for shape in iter_shapes(s.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        # 清理单独的 "D" 和 "esign" 碎片
        if text == "D" or text == "esign" or text == "esig":
            set_text_keep_fmt(shape, "")
    prs.save(PATH)

    # slide 4: 清理 3 处 "选题意义概述" 残留
    replace_sequential(4, ["选题意义概述"], [
        "设备自动触发交易，本地识别称重计价",
        "MQTT 上报，后台统一管理数据",
        "形成完整业务闭环",
    ])

    # slide 5: 清理 4 处 "选题意义概述" 残留
    replace_sequential(5, ["选题意义概述"], [
        "RK3576：边缘 AI 计算核心",
        "HX711：高精度称重采集",
        "SYN6288：语音播报与人机交互",
        "YOLO11：商品检测模型",
    ])

    # slide 6: 清理残留的 "决策力"、"远见力" SWOT 标签
    clean_slide(6, [
        (["决策力"], "触发"),
        (["远见力"], "播报"),
        (["影响力"], "计价"),
    ])

    # slide 9: 清理 "添加标题"（3处）和 "您的内容打在这里"（3处）
    replace_sequential(9, ["添加标题"], [
        "重量稳定后刷新摄像头缓冲帧并拍照",
        "查询单价并计算总价",
        "写入 JSONL 记录并通过 MQTT 上报",
    ])
    replace_sequential(9, ["您的内容打在这里"], [
        "空秤待机时预加载 RKNN runtime，保持摄像头常开",
        "重量超过阈值后多次采样确认稳定，避免误触发",
        "NPU 推理识别商品，HX711 读取重量参与计价",
    ])

    # slide 14: 清理 "添加标题"（1处）
    replace_sequential(14, ["添加标题"], [
        "后台改策略，板端下一笔交易即生效",
    ])

    # slide 15: 清理 "您的内容打在这里"
    clean_slide(15, [
        (["您的内容打在这里"], "交易、事件、语音命令用普通消息；参数、策略用 Retained 消息，新设备上线即取最新规则"),
    ])

    # slide 17: 清理 "案例四" 残留
    clean_slide(17, [
        (["案例四"], "设备状态"),
    ])

    # slide 23: 清理谢谢聆听页的英文模板文字
    clean_slide(23, [
        (["year-end summary", "About the summary", "boutique PPT"], "Smart-Cat 边缘智能称重识别一体机"),
    ])

    # 最终全局扫描：清理任何剩余的通用占位符
    prs = Presentation(PATH)
    placeholder_patterns = [
        "添加您的标题", "添加标题", "您的内容打在这里", "选题意义概述", "选题背景概述",
        "案例一", "案例二", "案例三", "案例四", "案例五",
        "文献一", "文献二", "文献三",
        "单击此处添加文字", "点击添加相关标题", "请替换文字内容",
    ]
    replacement_pool = [
        "Smart-Cat 边缘智能称重识别系统",
        "RK3576 NPU 边缘推理",
        "MQTT 云边协同",
        "Web 后台管理",
        "语音补盲与人工确认",
    ]
    pool_idx = 0
    for slide in prs.slides:
        for shape in iter_shapes(slide.shapes):
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            for pat in placeholder_patterns:
                if pat in text:
                    set_text_keep_fmt(shape, replacement_pool[pool_idx % len(replacement_pool)])
                    pool_idx += 1
                    break
    prs.save(PATH)

    print("最终清理完成")


if __name__ == "__main__":
    main()
