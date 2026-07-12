"""修正 PPT 中方案选型表的"推理位置"行。"""
from pptx import Presentation

SRC = r'd:\vscode\Smart-Cat\docs\Smart-Cat边缘智能称重识别一体机.补充版.pptx'
OUT = r'd:\vscode\Smart-Cat\docs\Smart-Cat边缘智能称重识别一体机.补充版.pptx'

OLD = "低延迟，弱网络可运行"
NEW = "低延迟，离线/弱网络下仍可独立运行"

p = Presentation(SRC)
changed = False
for i, slide in enumerate(p.slides, 1):
    if i != 10:    # 方案选型表在 slide 10
        continue
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        tbl = shape.table
        for row in tbl.rows:
            cells = [cell.text.strip() for cell in row.cells]
            # 第 2 行第 4 列 = "主要原因为低延迟..." 的位置
            if cells and cells[0] == "推理位置":
                target = row.cells[3]
                if OLD in target.text:
                    # 替换文字(保留原段落格式)
                    for para in target.text_frame.paragraphs:
                        for run in para.runs:
                            if OLD in run.text:
                                run.text = run.text.replace(OLD, NEW)
                                changed = True
                    if not changed:
                        # 如果文字在多 run 里被拆开,直接重写整段
                        target.text_frame.paragraphs[0].text = NEW
                        changed = True

if changed:
    p.save(OUT)
    print(f"已修正:  Slide 10, 推理位置行的主要原因")
    print(f"   旧: {OLD}")
    print(f"   新: {NEW}")
else:
    print("未找到匹配文字,请检查 PPT")
