"""校验 Smart-Cat演示.pptx 的每一页文本内容，标记空占位符或异常内容。"""
from pathlib import Path
from pptx import Presentation

PPT_PATH = Path(r"d:\vscode\Smart-Cat\docs\Smart-Cat演示.pptx")

# 可能残留的模板占位符关键词
PLACEHOLDER_KEYWORDS = [
    "添加您的标题", "您的内容打在这里", "案例一", "案例二", "案例三",
    "文献一", "文献二", "文献三", "选题意义概述", "本文论点", "本文结构",
    "本文结语", "优势", "劣势", "机会", "威胁", "点击", "占位", "placeholder",
    "Loram", "Ipsum", "Lorem",
]


def iter_shapes(shapes):
    """递归遍历所有形状，包括组合内的子形状。"""
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:  # GROUP
            yield from iter_shapes(shape.shapes)


def extract_text(shape):
    """安全提取形状文本。"""
    try:
        if shape.has_text_frame:
            return shape.text_frame.text.strip()
    except Exception:
        return ""
    return ""


def main():
    if not PPT_PATH.exists():
        print(f"文件不存在: {PPT_PATH}")
        return

    prs = Presentation(str(PPT_PATH))
    print(f"PPT 总页数: {len(prs.slides)}")
    print("=" * 80)

    for idx, slide in enumerate(prs.slides):
        print(f"\n--- 第 {idx} 页 ---")
        texts = []
        for shape in iter_shapes(slide.shapes):
            text = extract_text(shape)
            if text:
                texts.append(text)

        if not texts:
            print("  [警告] 本页无任何文本内容")
            continue

        for t in texts:
            # 检测占位符
            is_placeholder = any(kw.lower() in t.lower() for kw in PLACEHOLDER_KEYWORDS)
            tag = " [可能占位符]" if is_placeholder else ""
            # 截断显示长文本
            display = t if len(t) <= 100 else t[:97] + "..."
            print(f"  -{tag} {display}")


if __name__ == "__main__":
    main()
